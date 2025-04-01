import os
from typing import Dict, Any, List, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
from pathlib import Path
import json

from app.services.llm_service import generate_slide_content

# Template styling
TEMPLATES = {
    "academic": {
        "slide_width": Inches(13.333),
        "slide_height": Inches(7.5),
        "title_font": "Arial",
        "title_font_size": Pt(36),
        "body_font": "Arial",
        "body_font_size": Pt(24),
        "colors": {
            "title_background": RGBColor(35, 75, 120),
            "title_text": RGBColor(255, 255, 255),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(0, 0, 0),
            "accent1": RGBColor(0, 112, 192),
            "accent2": RGBColor(237, 125, 49)
        }
    },
    "minimalist": {
        "slide_width": Inches(13.333),
        "slide_height": Inches(7.5),
        "title_font": "Calibri",
        "title_font_size": Pt(40),
        "body_font": "Calibri",
        "body_font_size": Pt(28),
        "colors": {
            "title_background": RGBColor(255, 255, 255),
            "title_text": RGBColor(0, 0, 0),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(0, 0, 0),
            "accent1": RGBColor(180, 180, 180),
            "accent2": RGBColor(100, 100, 100)
        }
    },
    "corporate": {
        "slide_width": Inches(13.333),
        "slide_height": Inches(7.5),
        "title_font": "Calibri",
        "title_font_size": Pt(36),
        "body_font": "Calibri",
        "body_font_size": Pt(24),
        "colors": {
            "title_background": RGBColor(31, 73, 125),
            "title_text": RGBColor(255, 255, 255),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(0, 0, 0),
            "accent1": RGBColor(0, 112, 192),
            "accent2": RGBColor(255, 192, 0)
        }
    }
}

async def generate_presentation(
    paper_content: Dict[str, Any],
    output_dir: str,
    template: str = "academic",
    output_format: str = "pptx",
    include_images: bool = True,
    custom_options: Dict[str, Any] = None
) -> str:
    """
    Generate a presentation from paper content.
    
    Args:
        paper_content: Paper content dictionary
        output_dir: Directory to save the output
        template: Template to use
        output_format: Output format (pptx, pdf)
        include_images: Whether to include images from the paper
        custom_options: Additional custom options
        
    Returns:
        Path to the generated presentation file
    """
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Get template configuration (default to academic if not found)
    template_config = TEMPLATES.get(template, TEMPLATES["academic"])
    
    # Create a new presentation
    prs = Presentation()
    
    # Set slide dimensions
    prs.slide_width = template_config["slide_width"]
    prs.slide_height = template_config["slide_height"]
    
    # Get slide content using LLM
    max_slides = custom_options.get("max_slides", 10) if custom_options else 10
    focus_areas = custom_options.get("focus_areas", None) if custom_options else None
    
    slides_content = await generate_slide_content(
        paper_content=paper_content,
        max_slides=max_slides,
        focus_areas=focus_areas
    )
    
    # Create slides
    for slide_data in slides_content:
        # Determine layout
        layout = slide_data.get("layout", "content")
        
        if layout == "title":
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.get("title", "Presentation")
            
            # Set subtitle if available
            if slide.placeholders[1]:  # Subtitle placeholder
                subtitle = slide.placeholders[1]
                content_text = "\n".join(slide_data.get("content", []))
                subtitle.text = content_text
            
        else:
            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.get("title", "Slide")
            
            # Set content
            content = slide.placeholders[1]  # Content placeholder
            tf = content.text_frame
            
            for i, item in enumerate(slide_data.get("content", [])):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = item
                p.level = 0  # Top level bullet
                
                # Apply template styling
                p.font.name = template_config["body_font"]
                p.font.size = template_config["body_font_size"]
        
        # Add speaker notes if available
        notes = slide_data.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    
    # Include images if requested
    if include_images and paper_content.get("images"):
        # Create an image slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        title = slide.shapes.title
        title.text = "Figures from the Paper"
        
        # Calculate image positions
        image_files = paper_content.get("images", [])
        
        # Only include up to 4 images for simplicity
        image_files = image_files[:4]
        num_images = len(image_files)
        
        if num_images > 0:
            # Clear placeholder content
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Skip title
                    if shape != title:
                        shape.text = ""
            
            # Simple grid layout based on number of images
            paper_dir = Path(f"data/papers/{paper_content.get('paper_id')}")
            
            # Calculate image positions
            if num_images == 1:
                # Single image centered
                left = Inches(3)
                top = Inches(2)
                width = Inches(7)
                height = Inches(4)
                
                img_path = paper_dir / image_files[0]
                if img_path.exists():
                    slide.shapes.add_picture(str(img_path), left, top, width, height)
            
            elif num_images == 2:
                # Two images side by side
                for i, img_file in enumerate(image_files):
                    left = Inches(1.5 + i * 7)
                    top = Inches(2)
                    width = Inches(5)
                    height = Inches(4)
                    
                    img_path = paper_dir / img_file
                    if img_path.exists():
                        slide.shapes.add_picture(str(img_path), left, top, width, height)
            
            elif num_images == 3 or num_images == 4:
                # 2x2 grid
                for i, img_file in enumerate(image_files):
                    row = i // 2
                    col = i % 2
                    
                    left = Inches(1.5 + col * 7)
                    top = Inches(1.5 + row * 3.5)
                    width = Inches(5)
                    height = Inches(3)
                    
                    img_path = paper_dir / img_file
                    if img_path.exists():
                        slide.shapes.add_picture(str(img_path), left, top, width, height)
    
    # Save presentation
    output_filename = f"presentation.{output_format}"
    output_path = os.path.join(output_dir, output_filename)
    
    prs.save(output_path)
    
    # Convert to PDF if requested and output format is pdf
    if output_format.lower() == "pdf":
        # Note: Direct PDF conversion is not supported by python-pptx
        # In a real application, you might use a tool like LibreOffice or a conversion API
        pass
    
    return output_filename

def apply_template_styling(presentation: Presentation, template: str):
    """
    Apply template styling to a presentation.
    
    Args:
        presentation: Presentation object
        template: Template name
    """
    template_config = TEMPLATES.get(template, TEMPLATES["academic"])
    
    # Set slide dimensions
    presentation.slide_width = template_config["slide_width"]
    presentation.slide_height = template_config["slide_height"]
    
    # Apply styling to all slides
    for slide in presentation.slides:
        # Style title
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_frame = title_shape.text_frame
            title_frame.text = title_frame.text  # Preserve text
            
            for paragraph in title_frame.paragraphs:
                paragraph.font.name = template_config["title_font"]
                paragraph.font.size = template_config["title_font_size"]
                paragraph.font.color.rgb = template_config["colors"]["title_text"]
        
        # Style content
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                text_frame = shape.text_frame
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.name = template_config["body_font"]
                    paragraph.font.size = template_config["body_font_size"]
                    paragraph.font.color.rgb = template_config["colors"]["text"]

def extract_text_for_slides(paper_content: Dict[str, Any]) -> Dict[str, str]:
    """
    Extract relevant text sections from paper content for slides.
    
    Args:
        paper_content: Paper content dictionary
        
    Returns:
        Dictionary with sections of text for different slide types
    """
    full_text = paper_content.get("text", "")
    sections = paper_content.get("sections", [])
    
    # Initialize result dictionary
    result = {
        "title": "",
        "abstract": "",
        "introduction": "",
        "methodology": "",
        "results": "",
        "discussion": "",
        "conclusion": "",
        "references": ""
    }
    
    # Try to extract title from the start of the document
    title_match = re.search(r"^(.+?)(?:\n|$)", full_text.strip())
    if title_match:
        result["title"] = title_match.group(1)
    
    # Extract sections based on section headings if available
    if sections:
        for section in sections:
            heading = section.get("heading", "").lower()
            content = section.get("content", "")
            
            if "abstract" in heading:
                result["abstract"] = content
            elif any(term in heading for term in ["introduction", "background"]):
                result["introduction"] = content
            elif any(term in heading for term in ["method", "approach", "experiment"]):
                result["methodology"] = content
            elif any(term in heading for term in ["result", "finding", "observation"]):
                result["results"] = content
            elif "discussion" in heading:
                result["discussion"] = content
            elif any(term in heading for term in ["conclusion", "summary", "future work"]):
                result["conclusion"] = content
            elif any(term in heading for term in ["reference", "bibliography"]):
                result["references"] = content
    
    # If sections weren't found, try to extract using regex patterns
    if not sections or all(v == "" for v in result.values()):
        # Try to find abstract
        abstract_match = re.search(r"(?:abstract|summary)[:\s]+(.+?)(?=\n\n|\n[A-Z]|\n\d|\Z)", 
                                  full_text, re.IGNORECASE | re.DOTALL)
        if abstract_match:
            result["abstract"] = abstract_match.group(1).strip()
        
        # Try to find introduction
        intro_match = re.search(r"(?:introduction|background)[:\s]+(.+?)(?=\n\n\d|\n\n[A-Z]|\Z)", 
                               full_text, re.IGNORECASE | re.DOTALL)
        if intro_match:
            result["introduction"] = intro_match.group(1).strip()
        
        # Try to find methods/methodology
        method_match = re.search(r"(?:method|methodology|materials and methods|approach|experiment)[:\s]+(.+?)(?=\n\n\d|\n\n[A-Z]|\Z)", 
                                full_text, re.IGNORECASE | re.DOTALL)
        if method_match:
            result["methodology"] = method_match.group(1).strip()
        
        # Try to find results
        results_match = re.search(r"(?:results|findings|observations)[:\s]+(.+?)(?=\n\n\d|\n\n[A-Z]|\Z)", 
                                 full_text, re.IGNORECASE | re.DOTALL)
        if results_match:
            result["results"] = results_match.group(1).strip()
        
        # Try to find discussion
        discussion_match = re.search(r"(?:discussion)[:\s]+(.+?)(?=\n\n\d|\n\n[A-Z]|\Z)", 
                                    full_text, re.IGNORECASE | re.DOTALL)
        if discussion_match:
            result["discussion"] = discussion_match.group(1).strip()
        
        # Try to find conclusion
        conclusion_match = re.search(r"(?:conclusion|conclusions|summary|future work)[:\s]+(.+?)(?=\n\n\d|\n\n[A-Z]|\Z)", 
                                    full_text, re.IGNORECASE | re.DOTALL)
        if conclusion_match:
            result["conclusion"] = conclusion_match.group(1).strip()
    
    return result
